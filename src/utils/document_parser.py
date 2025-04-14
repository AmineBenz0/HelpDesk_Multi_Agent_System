# utils/document_parser.py
import os
from typing import List, Optional
from pathlib import Path
import chromadb
from chromadb.config import Settings
from langchain.text_splitter import RecursiveCharacterTextSplitter

class DocumentProcessor:
    def __init__(self, persist_dir: str = "chroma_db"):
        self.persist_dir = persist_dir
        self.chroma_client = chromadb.Client(
            Settings(persist_directory=persist_dir, chroma_db_impl="duckdb+parquet")
        )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )

    def process_document(self, file_path: str) -> List[str]:
        """Process a document based on its file type."""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            from PyPDF2 import PdfReader
            text = ""
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                text = "\n".join([page.extract_text() for page in reader.pages])
            return self._chunk_text(text)
            
        elif ext in ('.doc', '.docx'):
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return self._chunk_text(text)
            
        elif ext == '.pptx':
            from pptx import Presentation
            text = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return self._chunk_text("\n".join(text))
        
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks with overlapping windows."""
        return self.text_splitter.split_text(text)

    def store_embeddings(self, chunks: List[str], collection_name: str) -> None:
        """Store document chunks as embeddings in ChromaDB."""
        collection = self.chroma_client.create_collection(collection_name)
        ids = [str(i) for i in range(len(chunks))]
        collection.add(
            documents=chunks,
            ids=ids
        )
        self.chroma_client.persist()

    def query_collection(self, query: str, collection_name: str, n_results: int = 3) -> List[str]:
        """Query the vector database for similar documents."""
        collection = self.chroma_client.get_collection(collection_name)
        results = collection.query(
            query_texts=[query],
            n_results=n_results
        )
        return results['documents'][0]